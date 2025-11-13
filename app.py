import io
import re
from typing import List, Optional
from urllib.parse import urljoin

import requests
from bs4 import BeautifulSoup
from readability import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
import streamlit as st
from openai import OpenAI

APP_VERSION = "press2ppt v1.9 (image resize mode)"

# ========= 設定 =========
TEMPLATE_PATH = "templates/cuprum_template.pptx"
DEFAULT_FONTS = ["Meiryo", "Yu Gothic UI", "MS UI Gothic", "Calibri"]

# スタイル設定
TITLE_COLOR = RGBColor(255, 255, 255)
TITLE_SIZE_PT = 28
TITLE_BOLD = True

BODY_COLOR = RGBColor(0, 153, 153)
BODY_SIZE_PT = 24
BODY_BOLD = True

# ロゴ等の除外設定
IMG_EXCLUDE_RE = re.compile(
    r"(?:^|[-_/])(logo|favicon|sprite|badge|mark|header|footer|og_image|common/images/og_image\.png)\b",
    re.IGNORECASE,
)
EXACT_IMG_BLACKLIST = {
    "https://www.jx-nmm.com/common/images/og_image.png",
    "http://www.jx-nmm.com/common/images/og_image.png",
}

# ========= 追加: HEIC/AVIF対応 =========
HEIF_OK = False
AVIF_OK = False
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
    HEIF_OK = True
except Exception:
    pass

try:
    import pillow_avif
    AVIF_OK = True
except Exception:
    pass

# ========= 追加: 画像トリミングUI（streamlit-cropper） =========
CROPPER_OK = False
try:
    from streamlit_cropper import st_cropper
    CROPPER_OK = True
except Exception:
    CROPPER_OK = False


def _open_image_any(file_obj) -> Optional[Image.Image]:
    try:
        if hasattr(file_obj, "read"):
            data = file_obj.read()
            try:
                file_obj.seek(0)
            except Exception:
                pass
        elif isinstance(file_obj, (bytes, bytearray)):
            data = file_obj
        else:
            data = file_obj.getvalue()
        buf = io.BytesIO(data)
        img = Image.open(buf)
        if getattr(img, "is_animated", False):
            img.seek(0)
        return img.convert("RGB")
    except Exception:
        return None


# ========= OpenAI =========
def get_client(api_key: Optional[str]):
    try:
        return OpenAI(api_key=api_key) if api_key else OpenAI()
    except Exception:
        return None


# ========= HTML取得 =========
def fetch_html_public(url: str) -> str:
    r = requests.get(url, timeout=20, headers={"User-Agent": "Mozilla/5.0"})
    r.raise_for_status()
    return r.text


def _best_from_srcset(srcset: str) -> Optional[str]:
    try:
        parts = [p.strip() for p in srcset.split(",") if p.strip()]
        pairs = []
        for p in parts:
            s = p.split()
            if len(s) == 1:
                pairs.append((s[0], 0))
            else:
                w = 0
                try:
                    if s[1].endswith("w"):
                        w = int(s[1][:-1])
                except Exception:
                    pass
                pairs.append((s[0], w))
        pairs.sort(key=lambda x: x[1], reverse=True)
        return pairs[0][0] if pairs else None
    except Exception:
        return None


def parse_page(url: str) -> dict:
    html = fetch_html_public(url)
    return _parse_common(html, base_url=url)


def _parse_common(html: str, base_url: str = "") -> dict:
    doc = Document(html)
    title = (doc.short_title() or "").strip()
    if not title:
        try:
            head = BeautifulSoup(html, "lxml").find("head")
            if head:
                t = head.find("title")
                if t and t.get_text(strip=True):
                    title = t.get_text(strip=True)
        except Exception:
            pass

    main_html = doc.summary(html_partial=True)
    soup = BeautifulSoup(main_html, "lxml")

    ps = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
    text = " ".join(ps)

    def abs_url(u: str) -> str:
        if not base_url:
            return u
        if u and u.startswith("/"):
            return urljoin(base_url, u)
        return u

    urls: List[str] = []
    for img in soup.find_all("img"):
        cand = img.get("src") or img.get("data-src") or img.get("data-original")
        if not cand and img.get("srcset"):
            cand = _best_from_srcset(img.get("srcset"))
        if cand:
            urls.append(abs_url(cand))

    for pic in soup.find_all("picture"):
        for src in pic.find_all("source"):
            srcset = src.get("srcset")
            if srcset:
                cand = _best_from_srcset(srcset)
                if cand:
                    urls.append(abs_url(cand))

    try:
        head = BeautifulSoup(html, "lxml").find("head")
        if head:
            for prop in ["og:image", "twitter:image"]:
                tag = head.find("meta", property=prop) or head.find(
                    "meta", attrs={"name": prop}
                )
                if tag and tag.get("content"):
                    urls.append(abs_url(tag["content"]))
    except Exception:
        pass

    cleaned = []
    for u in urls:
        if not u or u.startswith("data:"):
            continue
        low = u.lower()
        u_noquery = re.sub(r"[?#].*$", "", u)
        if u_noquery in EXACT_IMG_BLACKLIST:
            continue
        if IMG_EXCLUDE_RE.search(low):
            continue
        cleaned.append(u)

    uniq, seen = [], set()
    for u in cleaned:
        base = re.sub(r"[?#].*$", "", u)
        if base not in seen:
            uniq.append(u)
            seen.add(base)

    return {"title": title, "text": text, "images": uniq}


# ========= 要約処理 =========
SYS_TITLER = (
    "あなたは日本語のPRアシスタントです。25文字を超える場合のみ自然な見出しに短縮。"
    "句読点含め25文字以内、固有名詞は優先して保持。"
)
SYS_SUMMARY = (
    "あなたは日本語のPRアシスタントです。企業プレスリリースの要旨を、"
    "指定の上限文字数以内で簡潔に要約してください。"
    "上限文字数の90%~95%の文字数に必ずして"
    "文を途中で切らず、句点「。」で完結させること。"
    "必要に応じて短文に分けても構いません。"
    "体言止めや重言は避け、固有名詞は保持。ですます調。"
    "JX金属株式会社やJX金属、JX等の主語は必ず省略し、それでも意味が通るように。"
    "冗長表現や重複を削り、意味を保ったまま上限以内に収めてください。"
)


def gpt_shorten_title(client: Optional[OpenAI], title: str) -> str:
    if len(title) <= 25 or not client:
        return title[:25]
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": SYS_TITLER},
                {"role": "user", "content": f"元タイトル：{title}\n25文字以内に短く"},
            ],
            temperature=0.2,
        )
        return (resp.choices[0].message.content or "").strip()[:25]
    except Exception:
        return title[:25]


def offline_summary(text: str) -> str:
    if not text:
        return ""
    sentences = re.split(r"(?<=[。！？!?])", text)
    chunk = "".join(sentences[:3])
    chunk = re.sub(r"\s+", " ", chunk)
    return chunk


def _tidy_clamp_to_limit(s: str, limit: int) -> str:
    s = s.strip()
    if len(s) <= limit:
        return s
    candidates = ["。", "！", "？", "!", "?", "…"]
    cut_pos = -1
    for ch in candidates:
        p = s.rfind(ch, 0, limit)
        if p > cut_pos:
            cut_pos = p
    if cut_pos >= 0 and cut_pos >= int(limit * 0.5):
        return s[:cut_pos + 1].strip()
    return s[:limit].rstrip("・、，,（(").rstrip()


def gpt_summarize_body(client: Optional[OpenAI], text: str, max_len: int = 120) -> str:
    head = (text or "")[:4000]
    base = offline_summary(head)
    if not client:
        return _tidy_clamp_to_limit(base, max_len)
    try:
        prompt_user = (
            f"{head}\n\n"
            f"上限{max_len}文字で、重要点を落とさず簡潔に要約してください。"
            f"文は途中で切らず、句点で完結させてください。"
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": SYS_SUMMARY},
                {"role": "user", "content": prompt_user},
            ],
            temperature=0.2,
        )
        s = (resp.choices[0].message.content or "").strip()
        return _tidy_clamp_to_limit(s, max_len)
    except Exception:
        return _tidy_clamp_to_limit(base, max_len)


def do_summary(text: str, max_len: int, api_key: Optional[str]) -> tuple[str, str]:
    client = get_client(api_key or None)
    if client:
        try:
            s = gpt_summarize_body(client, text, max_len)
            return s, "GPT"
        except Exception:
            pass
    s = gpt_summarize_body(None, text, max_len)
    return s, "OFFLINE"


# ========= 画像DL =========
def download_images(urls: List[str], limit: int = 4) -> List[Image.Image]:
    imgs: List[Image.Image] = []
    for u in urls[:limit]:
        try:
            r = requests.get(u, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            img = Image.open(io.BytesIO(r.content)).convert("RGB")
            if img.width < 300 or img.height < 180:
                continue
            imgs.append(img)
        except Exception:
            continue
    return imgs


# ========= PowerPoint生成 =========
def _first_placeholder(slide, types: tuple[int, ...]) -> Optional[object]:
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in types:
                return ph
        except Exception:
            continue
    return None


def _set_text(shape, text: str, size_pt: int, color: RGBColor, bold: bool):
    try:
        if not getattr(shape, "has_text_frame", False):
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        for fn in DEFAULT_FONTS:
            try:
                run.font.name = fn
                break
            except Exception:
                continue
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color
        run.font.bold = bold
    except Exception:
        pass


def get_layout_by_name(prs, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return None


def _place_image_contain(slide, ph, img: Image.Image):
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="JPEG", quality=92)
    buf.seek(0)

    frame_w = ph.width
    frame_h = ph.height
    frame_left = ph.left
    frame_top = ph.top

    iw, ih = img.size
    a_img = iw / ih
    a_frame = frame_w / frame_h

    if a_img >= a_frame:
        w = frame_w
        h = int(frame_w / a_img)
    else:
        h = frame_h
        w = int(frame_h * a_img)

    left = frame_left + int((frame_w - w) / 2)
    top = frame_top + int((frame_h - h) / 2)

    pic = slide.shapes.add_picture(buf, left, top, width=w, height=h)
    try:
        ph.element.getparent().remove(ph.element)
    except Exception:
        pass
    return pic


def build_pptx(
    template_path: str,
    title: str,
    summary: str,
    images: List[Image.Image],
    fit_mode: str,
) -> bytes:
    prs = Presentation(template_path)

    n = min(len(images), 3)
    layout_name = {
        0: "Cuprum Title+Body",
        1: "Cuprum Title+Body+1Pic",
        2: "Cuprum Title+Body+2Pic",
        3: "Cuprum Title+Body+3Pic",
    }[n]
    layout = get_layout_by_name(prs, layout_name) or prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_ph = _first_placeholder(slide, (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
    if title_ph is None:
        for ph in slide.placeholders:
            if "タイトル" in getattr(ph, "name", ""):
                title_ph = ph
                break
    if title_ph is not None:
        _set_text(title_ph, title, TITLE_SIZE_PT, TITLE_COLOR, TITLE_BOLD)

    body_ph = _first_placeholder(slide, (PP_PLACEHOLDER.BODY,))
    if body_ph is None:
        body_ph = _first_placeholder(slide, (PP_PLACEHOLDER.CONTENT,))
    if body_ph is None:
        for ph in slide.placeholders:
            if "テキスト" in getattr(ph, "name", ""):
                body_ph = ph
                break
    if body_ph is not None:
        _set_text(body_ph, summary, BODY_SIZE_PT, BODY_COLOR, BODY_BOLD)

    pic_placeholders = []
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                pic_placeholders.append(ph)
        except Exception:
            continue
    pic_placeholders.sort(key=lambda sh: (sh.left, sh.top))

    use_n = min(len(images), len(pic_placeholders))
    for i in range(use_n):
        ph = pic_placeholders[i]
        img = images[i]
        if fit_mode.startswith("収める"):
            _place_image_contain(slide, ph, img)
        else:
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="JPEG")
            buf.seek(0)
            try:
                ph.insert_picture(buf)
            except Exception:
                slide.shapes.add_picture(buf, ph.left, ph.top, width=ph.width)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()


# ========= 画像リサイズ用ヘルパー =========
def crop_to_16_9_center(img: Image.Image) -> Image.Image:
    """中央を16:9にトリミング（縦横どちらでもOK）"""
    w, h = img.size
    target_ratio = 16 / 9
    current_ratio = w / h

    if current_ratio > target_ratio:
        # 横に長い → 幅を削る
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        box = (left, 0, left + new_w, h)
    else:
        # 縦に長い → 高さを削る
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        box = (0, top, w, top + new_h)

    return img.crop(box)


def resize_long_side(img: Image.Image, long_side_px: int) -> Image.Image:
    """長辺を指定pxに合わせてリサイズ（アスペクト比維持）"""
    w, h = img.size
    long_side = max(w, h)
    if long_side <= 0:
        return img
    scale = long_side_px / long_side
    new_w = int(w * scale)
    new_h = int(h * scale)
    return img.resize((new_w, new_h), Image.LANCZOS)


# ========= UI =========
st.set_page_config(
    page_title="プレスURL / コピペ / 画像リサイズ → Cuprum PPT",
    page_icon="🧩",
    layout="wide",
)
st.title(f"プレスURL / コピペ＋画像 / リサイズ → Cuprumテンプレ自動作成｜{APP_VERSION}")

with st.sidebar:
    st.header("設定")
    template_file = st.file_uploader("テンプレ（.pptx）を差し替え可", type=["pptx"])
    api_key = st.text_input("OpenAI API Key（未入力/失敗時はオフライン要約）", type="password")
    max_images = st.slider("最大画像数（先頭から使用、上限3枚）", 0, 6, 3)
    summary_length = st.slider("要約文字数上限（目安）", 120, 400, 160, 20)
    fit_mode = st.selectbox(
        "画像のはめ込み方法",
        ["収める（余白あり・全体表示）", "埋める（トリミングあり）"],
        index=0,
    )
    show_debug = st.checkbox("🧩 デバッグ出力を表示", value=True)
    st.caption("タイトル>25文字は短縮。本文は上限文字数で要約（コピペ版は要約オプション）。")
    st.markdown("---")
    st.subheader("画像リサイズモード用情報")
    st.write(f"HEIC/HEIF対応: {'✅' if HEIF_OK else '⚠️ pillow-heif未導入'}")
    st.write(f"AVIF対応: {'✅' if AVIF_OK else '⚠️ pillow-avif未導入'}")
    st.write(f"トリミングUI: {'✅ streamlit-cropper有効' if CROPPER_OK else '⚠️ 自動センタークロップのみ'}")

mode = st.radio(
    "入力モード",
    ["プレスリリースモード", "Sharepointコピペモード", "画像リサイズモード"],
    horizontal=True,
)

# 共有の作業用変数（PPT用）
title_final = ""
summary_final = ""
engine_used = "NO_SUMMARY"
images: List[Image.Image] = []
parsed_preview = None

# ============== モード1：URLモード ==============
if mode == "プレスリリースモード":
    url = st.text_input("プレスリリースのURL（社外サイト推奨）")
    if st.button("① 内容を抽出（URLから）"):
        try:
            parsed = parse_page(url)
            st.session_state["parsed_url"] = parsed
            st.success("抽出しました。下で要約・画像を確認してください。")
        except Exception as e:
            st.error(f"抽出に失敗しました: {type(e).__name__}: {e}")

    parsed = st.session_state.get("parsed_url")
    if parsed:
        st.subheader("抽出結果プレビュー（URLモード）")
        left, right = st.columns([2, 1])
        with left:
            st.write("抽出タイトル:", parsed.get("title") or "(なし)")
            raw_text = parsed.get("text") or ""
            st.write(
                "本文（先頭プレビュー）:",
                raw_text[:300] + ("…" if len(raw_text) > 300 else ""),
            )
        with right:
            st.write("候補画像URL（先頭から使用）")
            candidates = parsed.get("images", [])
            if candidates:
                for i, u in enumerate(candidates[:max_images]):
                    st.write(f"{i+1}. {u}")
            else:
                st.write("（画像候補なし）")

        client = get_client(api_key or None)
        title_final = gpt_shorten_title(client, parsed.get("title") or "（無題）")

        summary_final, engine_used = do_summary(
            parsed.get("text") or "", summary_length, api_key
        )
        st.info(
            f"要約エンジン: {engine_used} / 原文: {len(parsed.get('text') or '')}文字 → 出力: {len(summary_final)}文字"
        )

        sel_urls = parsed.get("images", [])[:max_images]
        images = download_images(sel_urls, limit=max_images)
        if images:
            cols = st.columns(min(len(images), 3))
            for i, img in enumerate(images):
                with cols[i % len(cols)]:
                    st.image(img, caption=f"Image {i+1}", use_container_width=True)
        else:
            st.info("表示可能な画像が見つかりませんでした。")

        parsed_preview = {"title": title_final, "summary": summary_final}

# ============== モード2：コピペ＋画像アップロード ==============
elif mode == "Sharepointコピペモード":
    manual_title = st.text_input("記事タイトル（コピペ）")
    manual_body = st.text_area("記事本文（コピペ）", height=220)
    colA, colB = st.columns(2)
    with colA:
        do_summarize = st.checkbox(
            "本文を要約する（サイドバーの上限文字数を使用）", value=True
        )
    with colB:
        do_shorten_title = st.checkbox(
            "タイトルが25文字超なら短縮する", value=True
        )

    uploaded_files = st.file_uploader(
        "画像をアップロード（最大3枚まで）",
        type=[
            "jpg",
            "jpeg",
            "png",
            "webp",
            "bmp",
            "tiff",
            "tif",
            "gif",
            "heic",
            "heif",
            "avif",
        ],
        accept_multiple_files=True,
    )

    if st.button("① プレビュー生成（コピペ版）"):
        if not manual_title and not manual_body:
            st.warning("タイトルまたは本文のどちらかは入力してください。")
        else:
            client = get_client(api_key or None)
            if do_shorten_title:
                title_final = gpt_shorten_title(
                    client, manual_title or "（無題）"
                )
            else:
                title_final = (manual_title or "（無題）")[:100]

            if do_summarize:
                summary_final, engine_used = do_summary(
                    manual_body or "", summary_length, api_key
                )
            else:
                summary_final = (manual_body or "")[:summary_length]
                engine_used = "NO_SUMMARY"

            images = []
            if uploaded_files:
                for f in uploaded_files[:3]:
                    img = _open_image_any(f)
                    if img is None:
                        continue
                    if img.width < 300 or img.height < 180:
                        continue
                    images.append(img)

            st.session_state["manual_preview"] = {
                "title": title_final,
                "summary": summary_final,
                "images_len": len(images),
                "engine": engine_used,
                "raw_len": len(manual_body or ""),
                "out_len": len(summary_final or ""),
            }
            st.session_state["manual_images"] = images
            st.success("プレビューを作成しました。下で確認してください。")

    manual_prev = st.session_state.get("manual_preview")
    images = st.session_state.get("manual_images", [])
    if manual_prev:
        st.subheader("抽出結果プレビュー（コピペ版）")
        st.write("**タイトル（最終）**:", manual_prev["title"])
        st.write(
            f"**本文（{manual_prev['out_len']}文字 / エンジン: {manual_prev['engine']} / 原文{manual_prev['raw_len']}文字）**:"
        )
        st.write(manual_prev["summary"])
        if images:
            cols = st.columns(min(len(images), 3))
            for i, img in enumerate(images):
                with cols[i % len(cols)]:
                    st.image(img, caption=f"Uploaded {i+1}", use_container_width=True)
        else:
            st.info("画像は未アップロードです。")

        title_final = manual_prev["title"]
        summary_final = manual_prev["summary"]
        engine_used = manual_prev["engine"]
        parsed_preview = manual_prev

# ============== モード3：画像リサイズモード ==============
else:
    st.subheader("画像リサイズモード（16:9トリミング＋JPEG変換）")

    resize_px = st.radio(
        "長辺ピクセル数を選択",
        [1280, 1000, 600],
        index=0,
        horizontal=True,
    )

    uploaded_photos = st.file_uploader(
        "iPhoneなどで撮影した写真をアップロード（複数可）",
        type=[
            "jpg", "jpeg", "png", "webp", "bmp", "tiff", "tif", "gif",
            "heic", "heif", "avif",
        ],
        accept_multiple_files=True,
    )

    processed_images: List[tuple[str, Image.Image]] = []

    if uploaded_photos:
        if "crop_results" not in st.session_state:
            st.session_state["crop_results"] = {}
        crop_results = st.session_state["crop_results"]

        for idx, f in enumerate(uploaded_photos):
            img = _open_image_any(f)
            if img is None:
                continue

            w, h = img.size
            ratio = w / h if h != 0 else 0
            target_ratio = 16 / 9
            is_landscape = w >= h
            is_almost_16_9 = is_landscape and abs(ratio - target_ratio) < 0.02

            st.write(f"### 画像 {idx+1}: {getattr(f, 'name', f'image_{idx+1}')}")
            st.image(img, caption="元画像プレビュー", use_container_width=True)

            # ▼ プレビュー用クロップ画像
            if is_almost_16_9:
                st.info("ほぼ16:9の横画像のため、そのままリサイズ対象にします。")
                preview_crop = img
            else:
                if CROPPER_OK:
                    st.info(
                        "16:9以外または縦画像のため、16:9でトリミングしてください。"
                        "枠を調整し、『このトリミングを確定』を押してください。"
                    )
                    preview_crop = st_cropper(
                        img,
                        aspect_ratio=(16, 9),
                        box_color="#00FF00",
                        key=f"cropper_{idx}",
                        realtime_update=True,
                        return_type="image",
                    )
                else:
                    st.warning(
                        "streamlit-cropper未導入のため、自動で中央16:9にトリミングします。"
                    )
                    preview_crop = crop_to_16_9_center(img)

            # ▼ トリミング確定
            confirm_key = f"confirm_crop_{idx}"
            if st.button("このトリミングを確定", key=confirm_key):
                crop_results[idx] = preview_crop
                st.success("この画像のトリミングを確定しました。")

            used_crop = crop_results.get(idx, preview_crop)

            # ▼ リサイズ処理
            resized = resize_long_side(used_crop, resize_px)
            st.image(
                resized,
                caption=f"リサイズ後プレビュー（長辺 {resize_px}px / JPEG出力）",
                use_container_width=True,
            )

            base_name = getattr(f, "name", f"image_{idx+1}")
            base_name = re.sub(r"\.[^.]+$", "", base_name)
            out_name = f"{base_name}_{resize_px}px.jpg"
            processed_images.append((out_name, resized))

    # ▼ ZIPダウンロード
    if processed_images:
        import zipfile

        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for fname, im in processed_images:
                img_bytes = io.BytesIO()
                im.save(img_bytes, format="JPEG", quality=92)
                img_bytes.seek(0)
                zf.writestr(fname, img_bytes.read())
        zip_buf.seek(0)

        st.success(f"{len(processed_images)} ファイルを処理し、ZIPにまとめました。")
        st.download_button(
            "リサイズ画像ZIPをダウンロード",
            data=zip_buf,
            file_name=f"resized_images_{resize_px}px.zip",
            mime="application/zip",
        )
    else:
        st.caption("画像をアップロードすると、ここにリサイズ結果とZIPダウンロードボタンが表示されます。")

# ============== 共通：PPT生成（モード1&2のみ） ==============
if mode in ["プレスリリースモード", "Sharepointコピペモード"]:
    st.markdown("---")
    if st.button("② PPTを作成してダウンロード"):
        try:
            tpl_path = TEMPLATE_PATH
            if template_file is not None:
                tpl_path = "uploaded_template.pptx"
                with open(tpl_path, "wb") as f:
                    f.write(template_file.read())

            import os

            if not os.path.exists(tpl_path):
                st.error(f"テンプレが見つかりません: {tpl_path}")
            elif not parsed_preview:
                st.error("先に①でプレビューを作成してください。")
            else:
                if show_debug:
                    st.write(
                        {
                            "layout_candidates": [
                                "Cuprum Title+Body",
                                "Cuprum Title+Body+1Pic",
                                "Cuprum Title+Body+2Pic",
                                "Cuprum Title+Body+3Pic",
                            ],
                            "images_count": len(images or []),
                            "fit_mode": fit_mode,
                            "title_preview": (
                                title_final[:40]
                                + ("…" if len(title_final) > 40 else "")
                            ),
                            "summary_preview": (
                                summary_final[:60]
                                + ("…" if len(summary_final) > 60 else "")
                            ),
                            "engine_used": engine_used,
                        }
                    )

                ppt_bytes = build_pptx(
                    tpl_path,
                    title_final or "（無題）",
                    summary_final or "",
                    images or [],
                    fit_mode=fit_mode,
                )
                if not isinstance(ppt_bytes, (bytes, bytearray)) or len(ppt_bytes) == 0:
                    st.error("PPT生成に失敗しました（データ不正または空ファイル）。")
                else:
                    st.success("PPTを生成しました。")
                    st.download_button(
                        "ダウンロード",
                        data=ppt_bytes,
                        file_name="press_auto.pptx",
                        mime=(
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation"
                        ),
                    )
        except Exception as e:
            import traceback

            st.error(f"PPT生成に失敗しました: {type(e).__name__}: {e}")
            st.code("".join(traceback.format_exc()))
    else:
        st.caption("① 抽出/プレビュー → ② PPT作成 の順で操作してください。")
